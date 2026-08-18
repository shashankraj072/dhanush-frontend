from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AI-Powered Fitness Tracker"
subtitle.text = "Project Architecture & Technologies"

# Slide 1: Project Overview
bullet_slide_layout = prs.slide_layouts[1]
slide1 = prs.slides.add_slide(bullet_slide_layout)
shapes = slide1.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Project Overview"
tf = body_shape.text_frame
tf.text = "A full-stack web application designed for real-time workout tracking and personalized fitness planning."

p = tf.add_paragraph()
p.text = "Core Functionality: Utilizes the device's webcam to capture user movements and provide live posture feedback and automatic repetition counting."
p.level = 0

p = tf.add_paragraph()
p.text = "Personalization: Dynamically generates custom workout routines and nutritional meal plans based on individual user profiles and fitness goals."
p.level = 0

# Slide 2: Technology Stack
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes2 = slide2.shapes
title_shape2 = shapes2.title
body_shape2 = shapes2.placeholders[1]

title_shape2.text = "Technology Stack & Architecture"
tf2 = body_shape2.text_frame
tf2.text = "Frontend UI: React.js (Vite) and React Router for seamless SPA navigation."

p2 = tf2.add_paragraph()
p2.text = "Backend API: Python with Flask, designed to run as lightweight Serverless Functions or as a standalone server."
p2.level = 0

p2 = tf2.add_paragraph()
p2.text = "Communication: REST APIs. The frontend uses HTML5 <canvas> to capture video frames, converts them to Base64 images, and streams them to the backend."
p2.level = 0

p2 = tf2.add_paragraph()
p2.text = "Data Storage: In-memory state management (designed to plug into PostgreSQL or MongoDB)."
p2.level = 0

# Slide 3: Core Algorithms & AI
slide3 = prs.slides.add_slide(bullet_slide_layout)
shapes3 = slide3.shapes
title_shape3 = shapes3.title
body_shape3 = shapes3.placeholders[1]

title_shape3.text = "Computer Vision & Tracking Algorithms"
tf3 = body_shape3.text_frame
tf3.text = "Pose Estimation: Designed around MediaPipe BlazePose to extract 33 3D skeletal landmarks from raw video frames."

p3 = tf3.add_paragraph()
p3.text = "Posture Analysis Algorithm: Calculates joint angles using spatial coordinates to ensure the user maintains proper form, generating a live 'Accuracy Percentage'."
p3.level = 0

p3 = tf3.add_paragraph()
p3.text = "Heuristic State Machine (Rep Counting): Uses a deterministic algorithm tracking positional states ('up' and 'down') across angle thresholds to automatically count repetitions."
p3.level = 0

# Slide 4: Key Application Features
slide4 = prs.slides.add_slide(bullet_slide_layout)
shapes4 = slide4.shapes
title_shape4 = shapes4.title
body_shape4 = shapes4.placeholders[1]

title_shape4.text = "System Capabilities & Features"
tf4 = body_shape4.text_frame
tf4.text = "Real-Time Visual Feedback: Alerts users instantly if their posture breaks form (e.g., 'Keep back straight')."

p4 = tf4.add_paragraph()
p4.text = "Workout Progress Tracking: Logs completed workouts, tracks estimated calories burned over time, and stores recent history."
p4.level = 0

p4 = tf4.add_paragraph()
p4.text = "Client-Side Optimization: Frame sampling is capped at ~700ms intervals to keep the application lightweight, ensure smooth UI performance, and prevent network overload."
p4.level = 0

# Save presentation
prs.save('Fitness_Tracker_Presentation.pptx')
print("Successfully generated Fitness_Tracker_Presentation.pptx")
